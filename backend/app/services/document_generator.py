"""Unified document generator.

Produces PDF, DOCX, XLSX, PPTX, CSV, MD, HTML, TXT, and JSON output
from a simple structured spec or raw markdown/plain text.

All heavy libraries (reportlab, python-docx, openpyxl, python-pptx)
are already declared in backend/requirements.txt.

Public API:
    generate(spec: dict) -> tuple[bytes, str, str]
        Returns (file_bytes, media_type, suffix).

Supported `spec` fields (all optional unless noted):
    format    : str  REQUIRED  one of pdf|docx|xlsx|pptx|csv|md|html|txt|json
    title     : str            document title / first sheet name / deck title
    content   : str            primary body (markdown for pdf/docx/md/html,
                               plain text for txt, raw rows for csv via newlines)
    sections  : list[dict]     [{heading, body, bullets:[...]}, ...]
    tables    : list[dict]     [{name, headers:[...], rows:[[...],...]}, ...]
                               Used for xlsx (one sheet each) and csv (first only).
    slides    : list[dict]     [{title, bullets:[...], notes}] for pptx.
    data      : Any            free-form data for json export.
"""

from __future__ import annotations

import csv as _csv
import io
import json as _json
import re
from html import escape as _html_escape
from typing import Tuple


# ─────────────────────────────────────────────────────────────────────
# Format metadata
# ─────────────────────────────────────────────────────────────────────

MEDIA_TYPES = {
    "pdf":  "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "csv":  "text/csv",
    "md":   "text/markdown",
    "html": "text/html",
    "txt":  "text/plain",
    "json": "application/json",
}

SUPPORTED_FORMATS = tuple(MEDIA_TYPES.keys())


# ─────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────

def _spec_to_lines(spec: dict) -> list[str]:
    """Flatten a spec into a list of markdown-ish lines for text/PDF/DOCX rendering."""
    lines: list[str] = []
    title = spec.get("title")
    if title:
        lines.append(f"# {title}")
        lines.append("")
    content = spec.get("content")
    if content:
        for ln in str(content).splitlines():
            lines.append(ln)
        lines.append("")
    for section in spec.get("sections") or []:
        heading = section.get("heading")
        if heading:
            lines.append(f"## {heading}")
            lines.append("")
        body = section.get("body")
        if body:
            for ln in str(body).splitlines():
                lines.append(ln)
            lines.append("")
        for bullet in section.get("bullets") or []:
            lines.append(f"- {bullet}")
        if section.get("bullets"):
            lines.append("")
    # trim trailing blanks
    while lines and not lines[-1]:
        lines.pop()
    return lines


def _safe_filename(name: str | None, fallback: str = "document") -> str:
    base = (name or fallback).strip() or fallback
    cleaned = re.sub(r"[^a-zA-Z0-9_.-]+", "-", base).strip("-")
    return cleaned or fallback


# ─────────────────────────────────────────────────────────────────────
# Format implementations
# ─────────────────────────────────────────────────────────────────────

def _build_txt(spec: dict) -> bytes:
    lines = _spec_to_lines(spec)
    # strip markdown headers for plain text
    cleaned: list[str] = []
    for ln in lines:
        if ln.startswith("# "):
            cleaned.append(ln[2:])
            cleaned.append("=" * len(ln[2:]))
        elif ln.startswith("## "):
            cleaned.append(ln[3:])
            cleaned.append("-" * len(ln[3:]))
        else:
            cleaned.append(ln)
    return ("\n".join(cleaned) + "\n").encode("utf-8")


def _build_md(spec: dict) -> bytes:
    lines = _spec_to_lines(spec)
    return ("\n".join(lines) + "\n").encode("utf-8")


def _build_json(spec: dict) -> bytes:
    payload = spec.get("data")
    if payload is None:
        payload = {k: v for k, v in spec.items() if k != "format"}
    return _json.dumps(payload, indent=2, default=str).encode("utf-8")


def _build_html(spec: dict) -> bytes:
    lines = _spec_to_lines(spec)
    body_parts: list[str] = []
    in_list = False
    for ln in lines:
        if ln.startswith("# "):
            if in_list:
                body_parts.append("</ul>")
                in_list = False
            body_parts.append(f"<h1>{_html_escape(ln[2:])}</h1>")
        elif ln.startswith("## "):
            if in_list:
                body_parts.append("</ul>")
                in_list = False
            body_parts.append(f"<h2>{_html_escape(ln[3:])}</h2>")
        elif ln.startswith("### "):
            if in_list:
                body_parts.append("</ul>")
                in_list = False
            body_parts.append(f"<h3>{_html_escape(ln[4:])}</h3>")
        elif ln.startswith("- "):
            if not in_list:
                body_parts.append("<ul>")
                in_list = True
            body_parts.append(f"<li>{_html_escape(ln[2:])}</li>")
        elif not ln.strip():
            if in_list:
                body_parts.append("</ul>")
                in_list = False
            body_parts.append("")
        else:
            if in_list:
                body_parts.append("</ul>")
                in_list = False
            body_parts.append(f"<p>{_html_escape(ln)}</p>")
    if in_list:
        body_parts.append("</ul>")

    title = _html_escape(spec.get("title") or "Document")
    html = (
        "<!doctype html>\n<html><head><meta charset='utf-8'>"
        f"<title>{title}</title>"
        "<style>body{font-family:system-ui,-apple-system,Segoe UI,Roboto,sans-serif;"
        "max-width:780px;margin:2rem auto;padding:0 1rem;color:#111;line-height:1.55}"
        "h1{border-bottom:2px solid #444;padding-bottom:.3rem}"
        "h2{margin-top:1.6rem;color:#333}"
        "code{background:#f4f4f4;padding:.1em .35em;border-radius:3px}"
        "table{border-collapse:collapse;margin:1rem 0}"
        "td,th{border:1px solid #ccc;padding:.35rem .6rem;text-align:left}"
        "</style></head><body>\n"
        + "\n".join(body_parts)
        + "\n</body></html>\n"
    )
    return html.encode("utf-8")


def _build_csv(spec: dict) -> bytes:
    tables = spec.get("tables") or []
    buf = io.StringIO()
    writer = _csv.writer(buf)
    if tables:
        table = tables[0]
        headers = table.get("headers") or []
        if headers:
            writer.writerow(headers)
        for row in table.get("rows") or []:
            writer.writerow(row)
    else:
        # fall back to splitting content by newlines / commas as-is
        content = spec.get("content") or ""
        for ln in str(content).splitlines():
            writer.writerow([cell for cell in ln.split(",")])
    return buf.getvalue().encode("utf-8")


def _build_xlsx(spec: dict) -> bytes:
    from openpyxl import Workbook  # type: ignore

    wb = Workbook()
    # remove default sheet to control naming
    default_ws = wb.active
    wb.remove(default_ws)

    tables = spec.get("tables") or []
    if not tables:
        # synthesize a single sheet from content
        ws = wb.create_sheet(title=(spec.get("title") or "Sheet1")[:31] or "Sheet1")
        content = spec.get("content") or ""
        for row_idx, ln in enumerate(str(content).splitlines(), start=1):
            for col_idx, cell in enumerate(ln.split(","), start=1):
                ws.cell(row=row_idx, column=col_idx, value=cell.strip())
    else:
        for i, table in enumerate(tables):
            name = (table.get("name") or f"Sheet{i + 1}")[:31] or f"Sheet{i + 1}"
            ws = wb.create_sheet(title=name)
            headers = table.get("headers") or []
            if headers:
                ws.append(list(headers))
            for row in table.get("rows") or []:
                ws.append(list(row))

    if not wb.sheetnames:
        wb.create_sheet(title="Sheet1")

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _build_docx(spec: dict) -> bytes:
    from docx import Document  # type: ignore

    doc = Document()
    title = spec.get("title")
    if title:
        doc.add_heading(str(title), level=0)

    content = spec.get("content")
    if content:
        for ln in str(content).splitlines():
            if not ln.strip():
                doc.add_paragraph("")
            elif ln.startswith("# "):
                doc.add_heading(ln[2:], level=1)
            elif ln.startswith("## "):
                doc.add_heading(ln[3:], level=2)
            elif ln.startswith("### "):
                doc.add_heading(ln[4:], level=3)
            elif ln.startswith("- "):
                doc.add_paragraph(ln[2:], style="List Bullet")
            else:
                doc.add_paragraph(ln)

    for section in spec.get("sections") or []:
        heading = section.get("heading")
        if heading:
            doc.add_heading(str(heading), level=1)
        body = section.get("body")
        if body:
            for ln in str(body).splitlines():
                if ln.strip():
                    doc.add_paragraph(ln)
                else:
                    doc.add_paragraph("")
        for bullet in section.get("bullets") or []:
            doc.add_paragraph(str(bullet), style="List Bullet")

    for table in spec.get("tables") or []:
        headers = table.get("headers") or []
        rows = table.get("rows") or []
        if not headers and not rows:
            continue
        col_count = max(len(headers), max((len(r) for r in rows), default=0))
        if col_count == 0:
            continue
        t = doc.add_table(rows=1 + len(rows), cols=col_count)
        t.style = "Light Grid Accent 1"
        if headers:
            for j, h in enumerate(headers):
                t.rows[0].cells[j].text = str(h)
        for i, row in enumerate(rows, start=1):
            for j, cell in enumerate(row):
                if j < col_count:
                    t.rows[i].cells[j].text = "" if cell is None else str(cell)
        doc.add_paragraph("")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_pptx(spec: dict) -> bytes:
    from pptx import Presentation  # type: ignore

    prs = Presentation()

    # Title slide
    title = spec.get("title") or "Presentation"
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = str(title)
    if len(slide.placeholders) > 1 and spec.get("content"):
        slide.placeholders[1].text = str(spec["content"]).splitlines()[0][:200]

    bullet_layout = prs.slide_layouts[1]
    slides = spec.get("slides") or []
    if not slides:
        # build slides from sections
        for section in spec.get("sections") or []:
            slides.append({
                "title": section.get("heading"),
                "bullets": section.get("bullets") or (
                    [section["body"]] if section.get("body") else []
                ),
                "notes": None,
            })

    for s in slides:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = str(s.get("title") or "")
        body_ph = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body_ph is not None:
            tf = body_ph.text_frame
            bullets = s.get("bullets") or []
            if bullets:
                tf.text = str(bullets[0])
                for b in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = str(b)
                    p.level = 0
        notes = s.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_pdf(spec: dict) -> bytes:
    from reportlab.lib.pagesizes import LETTER  # type: ignore
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle  # type: ignore
    from reportlab.lib.units import inch  # type: ignore
    from reportlab.platypus import (  # type: ignore
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    )
    from reportlab.lib import colors  # type: ignore
    from xml.sax.saxutils import escape as xml_escape

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=LETTER,
        leftMargin=0.7 * inch, rightMargin=0.7 * inch,
        topMargin=0.6 * inch, bottomMargin=0.6 * inch,
        title=spec.get("title") or "Document",
    )
    styles = getSampleStyleSheet()
    h0 = ParagraphStyle("h0", parent=styles["Title"], fontSize=20, leading=24, spaceAfter=12)
    h1 = ParagraphStyle("h1", parent=styles["Heading1"], fontSize=14, leading=18, spaceAfter=6)
    h2 = ParagraphStyle("h2", parent=styles["Heading2"], fontSize=12, leading=15, spaceAfter=4)
    body = ParagraphStyle("body", parent=styles["BodyText"], fontSize=10.5, leading=14)
    bullet = ParagraphStyle("bullet", parent=body, leftIndent=14, bulletIndent=2)

    story: list = []
    title = spec.get("title")
    if title:
        story.append(Paragraph(xml_escape(str(title)), h0))

    def _render_line(ln: str):
        if not ln.strip():
            story.append(Spacer(1, 0.08 * inch))
        elif ln.startswith("# "):
            story.append(Paragraph(xml_escape(ln[2:]), h1))
        elif ln.startswith("## "):
            story.append(Paragraph(xml_escape(ln[3:]), h2))
        elif ln.startswith("### "):
            story.append(Paragraph(f"<b>{xml_escape(ln[4:])}</b>", body))
        elif ln.startswith("- "):
            story.append(Paragraph(f"&#8226; {xml_escape(ln[2:])}", bullet))
        else:
            story.append(Paragraph(xml_escape(ln), body))

    content = spec.get("content")
    if content:
        for ln in str(content).splitlines():
            _render_line(ln)

    for section in spec.get("sections") or []:
        heading = section.get("heading")
        if heading:
            story.append(Spacer(1, 0.08 * inch))
            story.append(Paragraph(xml_escape(str(heading)), h1))
        sec_body = section.get("body")
        if sec_body:
            for ln in str(sec_body).splitlines():
                _render_line(ln)
        for b in section.get("bullets") or []:
            story.append(Paragraph(f"&#8226; {xml_escape(str(b))}", bullet))

    for table in spec.get("tables") or []:
        headers = table.get("headers") or []
        rows = table.get("rows") or []
        if not headers and not rows:
            continue
        data: list[list[str]] = []
        if headers:
            data.append([str(h) for h in headers])
        for row in rows:
            data.append(["" if c is None else str(c) for c in row])
        if not data:
            continue
        story.append(Spacer(1, 0.12 * inch))
        tbl = Table(data, repeatRows=1 if headers else 0)
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]))
        story.append(tbl)

    if not story:
        story.append(Paragraph("(empty document)", body))

    doc.build(story)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────
# Public entry point
# ─────────────────────────────────────────────────────────────────────

_BUILDERS = {
    "pdf":  _build_pdf,
    "docx": _build_docx,
    "xlsx": _build_xlsx,
    "pptx": _build_pptx,
    "csv":  _build_csv,
    "md":   _build_md,
    "html": _build_html,
    "txt":  _build_txt,
    "json": _build_json,
}


def generate(spec: dict) -> Tuple[bytes, str, str]:
    """Generate a document. Returns (bytes, media_type, suffix)."""
    fmt = (spec.get("format") or "").strip().lower()
    if fmt not in _BUILDERS:
        raise ValueError(
            f"Unsupported format '{fmt}'. Supported: {', '.join(SUPPORTED_FORMATS)}"
        )
    builder = _BUILDERS[fmt]
    content = builder(spec)
    return content, MEDIA_TYPES[fmt], fmt


def filename_for(spec: dict, suffix: str) -> str:
    base = _safe_filename(spec.get("filename") or spec.get("title"), fallback="document")
    if base.lower().endswith("." + suffix):
        return base
    return f"{base}.{suffix}"


def render_conversation(
    title: str,
    messages: list[dict],
    fmt: str,
) -> Tuple[bytes, str, str, str]:
    """Render a chat conversation into the requested format.

    `messages` is a list of {role, content, created_at?} dicts.
    Returns (bytes, media_type, suffix, filename).
    """
    lines: list[str] = []
    if title:
        lines.append(f"# {title}")
        lines.append("")
    for m in messages:
        role = m.get("role") or "user"
        speaker = "You" if role == "user" else ("AI" if role == "assistant" else role.capitalize())
        ts = m.get("created_at") or ""
        header = f"## {speaker}" + (f"  _(at {ts})_" if ts else "")
        lines.append(header)
        lines.append("")
        body = (m.get("content") or "").rstrip()
        if body:
            for ln in body.splitlines():
                lines.append(ln)
        lines.append("")

    spec = {
        "format": fmt,
        "title": title or "Conversation",
        "content": "\n".join(lines),
        "filename": _safe_filename(title or "conversation", fallback="conversation"),
    }
    data, media_type, suffix = generate(spec)
    return data, media_type, suffix, filename_for(spec, suffix)
